from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer, AutoConfig
import torch
import json
import logging
import os
from werkzeug.utils import secure_filename
from datetime import datetime

import pdfplumber
import docx
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import bcrypt

# Import configuration
from config import (
    BASE_DIR, FLASK_DEBUG, HOST, PORT,
    SUMMARY_MODEL_PATH, QUESTION_GENERATION_MODEL_PATH, QUIZ_MODEL_PATH,
    HUGGINGFACE_SUMMARY_MODEL, HUGGINGFACE_QUESTION_MODEL, HUGGINGFACE_QUIZ_MODEL,
    USE_LOCAL_MODELS, FALLBACK_TO_HUGGINGFACE, model_exists, get_torch_dtype,
    UPLOAD_FOLDER, STUDY_MATERIALS_FOLDER, ALLOWED_EXTENSIONS
)

from db import users_collection, teachers_collection, materials_collection, db

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True, allow_headers="*", methods=["GET", "POST", "OPTIONS"])
logging.basicConfig(level=logging.DEBUG)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Log configuration
app.logger.info("=== Study Buddy Configuration ===")
app.logger.info(f"Use Local Models: {USE_LOCAL_MODELS}")
app.logger.info(f"Fallback to HuggingFace: {FALLBACK_TO_HUGGINGFACE}")
app.logger.info(f"Summary Model Path: {SUMMARY_MODEL_PATH}")
app.logger.info(f"Question Model Path: {QUESTION_GENERATION_MODEL_PATH}")
app.logger.info(f"Quiz Model Path: {QUIZ_MODEL_PATH}")

# Global model variables (will be loaded on first use)
summary_model = None
summary_tokenizer = None
question_model = None
question_tokenizer = None
quiz_model = None
quiz_tokenizer = None

def load_model_with_fallback(model_path: str, huggingface_model: str, model_name: str):
    """
    Load a model with fallback to HuggingFace if local model fails.
    
    Args:
        model_path: Path to local model
        huggingface_model: HuggingFace model identifier
        model_name: Name of the model (for logging)
    
    Returns:
        tuple: (model, tokenizer)
    """
    model = None
    tokenizer = None
    
    # Try to load local model first if enabled
    if USE_LOCAL_MODELS and model_exists(model_path):
        try:
            app.logger.info(f"Attempting to load local {model_name} model from {model_path}")
            
            # Load tokenizer
            tokenizer = AutoTokenizer.from_pretrained(model_path, trust_remote_code=True)
            
            # Try to load model
            try:
                model = AutoModelForSeq2SeqLM.from_pretrained(
                    model_path,
                        trust_remote_code=True,
                    torch_dtype=get_torch_dtype(),
                    ignore_mismatched_sizes=True
                    )
                app.logger.info(f"Successfully loaded local {model_name} model")
            except Exception as e1:
                app.logger.warning(f"Failed to load local model with default settings: {e1}")
                # Try with config first
                try:
                    config = AutoConfig.from_pretrained(model_path, trust_remote_code=True)
                    model = AutoModelForSeq2SeqLM.from_pretrained(
                        model_path,
                        config=config,
                        trust_remote_code=True,
                        torch_dtype=get_torch_dtype(),
                        ignore_mismatched_sizes=True
                    )
                    app.logger.info(f"Successfully loaded local {model_name} model with config")
                except Exception as e2:
                    app.logger.warning(f"Failed to load local model with config: {e2}")
                    model = None
        except Exception as e:
            app.logger.warning(f"Error loading local {model_name} model: {e}")
            model = None
            tokenizer = None
    
    # Fallback to HuggingFace model
    if (model is None or tokenizer is None) and FALLBACK_TO_HUGGINGFACE:
        try:
            app.logger.info(f"Loading HuggingFace {model_name} model: {huggingface_model}")
            tokenizer = AutoTokenizer.from_pretrained(huggingface_model)
            model = AutoModelForSeq2SeqLM.from_pretrained(huggingface_model)
            app.logger.info(f"Successfully loaded HuggingFace {model_name} model: {huggingface_model}")
        except Exception as e:
            app.logger.error(f"Failed to load HuggingFace {model_name} model: {e}")
            raise Exception(f"Failed to load {model_name} model from both local and HuggingFace sources. Error: {e}")
    
    if model is None:
        raise Exception(f"Failed to load {model_name} model. Check configuration.")
    
    model.eval()
    return model, tokenizer

def load_summary_model():
    """Load the fine-tuned summary model with fallback"""
    global summary_model, summary_tokenizer
    if summary_model is None:
        try:
            summary_model, summary_tokenizer = load_model_with_fallback(
                        SUMMARY_MODEL_PATH,
                HUGGINGFACE_SUMMARY_MODEL,
                "summary"
            )
            app.logger.info("Summary model ready for use")
        except Exception as e:
            app.logger.error(f"Error loading summary model: {e}")
            raise

def load_question_model():
    """Load the fine-tuned question generation model with fallback"""
    global question_model, question_tokenizer
    if question_model is None:
        try:
            question_model, question_tokenizer = load_model_with_fallback(
                        QUESTION_GENERATION_MODEL_PATH,
                HUGGINGFACE_QUESTION_MODEL,
                "question generation"
            )
            app.logger.info("Question generation model ready for use")
        except Exception as e:
            app.logger.error(f"Error loading question generation model: {e}")
            raise

def load_quiz_model():
    """Load the fine-tuned quiz generation model with fallback"""
    global quiz_model, quiz_tokenizer
    if quiz_model is None:
        try:
            quiz_model, quiz_tokenizer = load_model_with_fallback(
                        QUIZ_MODEL_PATH,
                HUGGINGFACE_QUIZ_MODEL,
                "quiz"
            )
            app.logger.info("Quiz model ready for use")
        except Exception as e:
            app.logger.error(f"Error loading quiz model: {e}")
            raise

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(filepath, start_page=1, end_page=None):
    text = ""
    try:
        if not os.path.exists(filepath):
            raise ValueError(f"File not found: {filepath}")
            
        if not os.access(filepath, os.R_OK):
            raise ValueError(f"File not readable: {filepath}")
            
        with pdfplumber.open(filepath) as pdf:
            total_pages = len(pdf.pages)
            app.logger.debug(f"Processing PDF with {total_pages} pages")
            
            if start_page < 1 or start_page > total_pages:
                raise ValueError(f"Start page must be between 1 and {total_pages}")
                
            if end_page is None:
                end_page = total_pages
            elif end_page < start_page or end_page > total_pages:
                raise ValueError(f"End page must be between {start_page} and {total_pages}")
                
            for page_num in range(start_page - 1, end_page):
                try:
                    page = pdf.pages[page_num]
                    page_text = page.extract_text()
                    
                    if page_text:
                        text += page_text + "\n"
                        app.logger.debug(f"Extracted {len(page_text)} characters from page {page_num + 1}")
                    else:
                        app.logger.warning(f"No text found on page {page_num + 1}")
                        
                except Exception as e:
                    app.logger.error(f"Error extracting text from page {page_num + 1}: {str(e)}")
                    continue
                    
        if not text.strip():
            raise ValueError("No extractable text found in the specified page range")
            
        app.logger.info(f"Successfully extracted text from pages {start_page} to {end_page}")
        return text.strip()
        
    except pdfplumber.PDFSyntaxError as e:
        app.logger.error(f"Invalid PDF file: {str(e)}")
        raise ValueError("The file appears to be corrupted or not a valid PDF")
    except Exception as e:
        app.logger.error(f"Error processing PDF: {str(e)}")
        raise ValueError(f"Failed to process PDF: {str(e)}")

def extract_text_from_docx(filepath):
    doc = docx.Document(filepath)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text.strip()

def extract_text_from_pptx(filepath, start_slide=1, end_slide=None):
    prs = Presentation(filepath)
    text = ""
    total_slides = len(prs.slides)
    if end_slide is None or end_slide > total_slides:
        end_slide = total_slides
    for slide_num in range(start_slide - 1, end_slide):
        slide = prs.slides[slide_num]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_text_from_url(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        paragraphs = soup.find_all("p")
        text = "\n".join([para.get_text() for para in paragraphs])
        return text.strip()
    except requests.exceptions.RequestException:
        return None

def extract_text_only(filepath, start_page=1, end_page=None):
    """Extract raw text from file without summarization"""
    file_ext = os.path.splitext(filepath)[1].lower()
    
    if file_ext == '.pdf':
        return extract_text_from_pdf(filepath, start_page, end_page)
    elif file_ext == '.docx':
        return extract_text_from_docx(filepath)
    elif file_ext == '.pptx':
        return extract_text_from_pptx(filepath, start_page, end_page)
    elif file_ext == '.txt':
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        raise ValueError(f"Unsupported file type: {file_ext}")

def generate_summary(text):
    """Generate summary using fine-tuned model"""
    if not text or not isinstance(text, str) or text.strip() == "":
        return "Invalid input. Please enter text."
    
    try:
        # Load model if not already loaded
        load_summary_model()
        
        # Preprocess text - remove extra whitespace and ensure minimum length
        text = ' '.join(text.split())
        if len(text) < 50:
            return "The text is too short to summarize effectively."
        
        # Tokenize input
        inputs = summary_tokenizer(
            text,
            return_tensors="pt",
            max_length=512,
            truncation=True,
            padding=True
        )
        
        # Generate summary
        with torch.no_grad():
            outputs = summary_model.generate(
                inputs["input_ids"],
                max_length=150,
                min_length=30,
                num_beams=4,
                early_stopping=True,
                do_sample=False
            )
        
        # Decode generated summary
        summary = summary_tokenizer.decode(outputs[0], skip_special_tokens=True)
        
        if not summary or summary.strip() == "":
            return "Unable to generate summary from the provided text."
            
        return summary.strip()
    except Exception as e:
        app.logger.error(f"Error during summarization: {e}")
        return f"An error occurred during summarization: {str(e)}"

def parse_mcq(generated_text):
    """Parse the generated text into question, options, and the correct answer."""
    lines = generated_text.split("\n")
    question = lines[0].replace("Question:", "").strip()
    options = [line.replace("Option:", "").strip() for line in lines[1:-1]]
    correct_answer = lines[-1].replace("Answer:", "").strip()
    return question, options, correct_answer

@app.route('/extract_text', methods=['POST'])
def extract_text():
    """Extract raw text from file without summarization"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        # Extract text based on file type
        file_ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
        start_page = request.form.get('start_page', type=int) or 1
        end_page = request.form.get('end_page', type=int) or None
        
        try:
            text = extract_text_only(filepath, start_page, end_page)
        finally:
            # Clean up uploaded file
            try:
                os.remove(filepath)
            except:
                pass
        
        return jsonify({'text': text, 'filename': filename})
        
    except Exception as e:
        app.logger.error(f"Error in extract_text endpoint: {str(e)}")
        return jsonify({'error': f'Error extracting text: {str(e)}'}), 500

@app.route('/summarize', methods=['POST'])
def summarize():
    """Summarize text from file, URL, or direct text input"""
    try:
        user_email = None
        if request.is_json:
            user_email = request.json.get('user_email')
        
        # Handle file upload
        if 'file' in request.files:
            file = request.files['file']
            if file.filename == '':
                return jsonify({'error': 'No file selected'}), 400
            
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            
            # Extract text based on file type
            file_ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
            start_page = request.form.get('start_page', type=int) or 1
            end_page = request.form.get('end_page', type=int) or None
            
            if file_ext == 'pdf':
                text = extract_text_from_pdf(filepath, start_page, end_page)
            elif file_ext == 'docx':
                text = extract_text_from_docx(filepath)
            elif file_ext == 'pptx':
                text = extract_text_from_pptx(filepath, start_page, end_page)
            elif file_ext == 'txt':
                with open(filepath, 'r', encoding='utf-8') as f:
                    text = f.read()
            else:
                return jsonify({'error': 'Unsupported file type'}), 400
            
            # Generate summary
            summary = generate_summary(text)
            
            # Save to database
            try:
                db.save_summary(
                    user_email=user_email,
                    source_type='file',
                    source_content=None,
                    source_file_path=filepath,
                    summary_text=summary
                )
            except Exception as db_error:
                app.logger.warning(f"Failed to save summary to database: {db_error}")
            
            # Clean up uploaded file
            try:
                os.remove(filepath)
            except:
                pass
            
            return jsonify({'summary': summary})
        
        # Handle URL
        if request.is_json and 'url' in request.json:
            url = request.json['url']
            text = extract_text_from_url(url)
            if not text:
                return jsonify({'error': 'Failed to extract text from URL'}), 400
            
            summary = generate_summary(text)
            
            # Save to database
            try:
                db.save_summary(
                    user_email=user_email,
                    source_type='url',
                    source_content=url,
                    source_file_path=None,
                    summary_text=summary
                )
            except Exception as db_error:
                app.logger.warning(f"Failed to save summary to database: {db_error}")
            
            return jsonify({'summary': summary})
        
        # Handle text input
        if request.is_json and 'text' in request.json:
            text = request.json['text']
            summary = generate_summary(text)
            
            # Save to database
            try:
                db.save_summary(
                    user_email=user_email,
                    source_type='text',
                    source_content=text[:500],  # Store first 500 chars
                    source_file_path=None,
                    summary_text=summary
                )
            except Exception as db_error:
                app.logger.warning(f"Failed to save summary to database: {db_error}")
            
            return jsonify({'summary': summary})
            
        return jsonify({'error': 'No text, file, or URL provided'}), 400
    except Exception as e:
        app.logger.error(f"Error in summarize endpoint: {str(e)}")
        db.log_operation(
            endpoint='/summarize',
            operation_type='summarize',
            user_email=None,
            details=str(e),
            status='error',
            error_message=str(e)
        )
        return jsonify({'error': f'Error processing request: {str(e)}'}), 500

@app.route('/quiz', methods=['POST'])
def generate_quiz():
    """Generate multiple MCQs using fine-tuned quiz generation model (fixed multi-output)."""
    try:
        data = request.get_json() if request.is_json else {}
        quiz_data = data.get('quiz_data') or data.get('text', '')
        num_questions = int(data.get('num_questions', 3))

        if not quiz_data.strip():
            return jsonify({'error': 'Context is required. Send {"quiz_data": "text"} or {"text": "text"}'}), 400

        # Load model if not already loaded
        load_quiz_model()

        # Prepare input exactly like training
        input_text = f"Generate {num_questions} multiple-choice questions from the following passage:\n{quiz_data}"

        # Tokenize input
        inputs = quiz_tokenizer.encode(
            input_text,
            return_tensors="pt",
            max_length=512,
            truncation=True,
            padding='max_length'
        )

        with torch.no_grad():
            # ✅ Use sampling (not beam search) for diverse multiple outputs
            outputs = quiz_model.generate(
                inputs,
                max_length=256,
                num_return_sequences=num_questions,
                do_sample=True,
                top_k=50,
                top_p=0.92,
                temperature=0.8,
                repetition_penalty=1.15,
                no_repeat_ngram_size=3,
                early_stopping=True
            )

        # Decode all generated MCQs separately
        generated_texts = [quiz_tokenizer.decode(o, skip_special_tokens=True) for o in outputs]

        quizzes = []
        for text in generated_texts:
            try:
                question, options, correct_answer = parse_mcq(text)
                quizzes.append({
                    'question': question,
                    'options': options,
                    'correct_answer': correct_answer,
                    'raw_output': text
                })
            except Exception:
                quizzes.append({
                    'question': text,
                    'options': ['Option A', 'Option B', 'Option C', 'Option D'],
                    'correct_answer': 'See raw output',
                    'raw_output': text
                })

        return jsonify({'quizzes': quizzes, 'quiz_list': quizzes})

    except Exception as e:
        app.logger.error(f"Error in quiz endpoint: {str(e)}")
        return jsonify({'error': f'Error generating quiz: {str(e)}'}), 500

@app.route('/question_answer', methods=['POST'])
def generate_question_answer():
    """Generate questions or answers using the fine-tuned T5 model."""
    try:
        # Parse JSON data
        data = request.get_json() if request.is_json else {}
        context = data.get('context') or data.get('text', '')
        question = data.get('question', '')
        num_questions = int(data.get('num_questions', 1))

        if not context and not question:
            return jsonify({'error': 'Either context or question must be provided'}), 400

        # Load model and tokenizer (if not loaded already)
        load_question_model()

        # ---- INPUT CONSTRUCTION ----
        if question:
            # ✅ Answer generation mode
            input_text = f"Context: {context}\nQuestion: {question}\nAnswer:"
        else:
            # ✅ Question generation mode — improved prompt
            input_text = (
                f"Task: question generation.\n"
                f"Generate {num_questions} meaningful and diverse questions "
                f"from the following passage:\n{context}"
            )

        # ---- TOKENIZATION ----
        input_ids = question_tokenizer.encode(
            input_text,
            return_tensors="pt",
            max_length=512,
            truncation=True,
            padding='max_length'
        )

        # ---- GENERATION ----
        with torch.no_grad():
            if not question:  # question generation
                outputs = question_model.generate(
                    input_ids,
                    max_length=512,
                    num_return_sequences=num_questions,
                    num_beams=4,
                    do_sample=True,           # add randomness for diversity
                    temperature=0.9,           # slightly higher temp for creativity
                    top_k=50,                  # sample from top 50 tokens
                    top_p=0.95,                # nucleus sampling
                    repetition_penalty=1.3,    # reduce repetition
                    no_repeat_ngram_size=2,    # avoid phrase repetition
                    early_stopping=True
                )
                # Decode all generated questions
                questions = [
                    question_tokenizer.decode(output, skip_special_tokens=True).strip()
                    for output in outputs
                ]

                # Filter duplicates
                unique_questions = list(dict.fromkeys(questions))

                return jsonify({
                    'questions': unique_questions,
                    'question_list': unique_questions
                })
            else:
                # ---- Answer generation ----
                outputs = question_model.generate(
                    input_ids,
                    max_length=256,
                    num_beams=5,
                    temperature=0.7,
                    top_k=50,
                    top_p=0.9,
                    repetition_penalty=1.2,
                    early_stopping=True
                )

                generated_text = question_tokenizer.decode(outputs[0], skip_special_tokens=True)
                answer = generated_text.replace(input_text, "").strip()

                return jsonify({'answer': answer})

    except Exception as e:
        app.logger.error(f"Error in question_answer endpoint: {str(e)}")
        return jsonify({'error': f'Error generating question/answer: {str(e)}'}), 500


@app.route('/important_questions', methods=['POST'])
def generate_important_questions():
    """Generate important questions using fine-tuned T5 model"""
    try:
        if not request.is_json:
            return jsonify({'error': 'Request must be JSON'}), 400
        
        data = request.get_json()
        important_data = data.get('important_data', '')
        
        if not important_data:
            return jsonify({'error': 'No important data provided'}), 400
        
        load_question_model()

        # Use the same prompt as used during training
        input_text = f"Generate questions from the following passage: {important_data}"
        
        # Tokenize input
        inputs = question_tokenizer(
            input_text,
            return_tensors="pt",
            max_length=512,
            truncation=True
        )

        # Generate multiple questions (same params as Streamlit)
        with torch.no_grad():
            outputs = question_model.generate(
                inputs["input_ids"],
                max_length=512,
                num_beams=5,
                num_return_sequences=3,
                early_stopping=True
            )

        # Decode all generated questions
        generated_questions = [
            question_tokenizer.decode(output, skip_special_tokens=True)
            for output in outputs
        ]
        
        return jsonify({'important_questions': generated_questions})

    except Exception as e:
        app.logger.error(f"Error in important_questions endpoint: {str(e)}")
        return jsonify({'error': f'Error generating important questions: {str(e)}'}), 500


@app.route('/signup', methods=['POST'])
def signup():
    """Register a new user"""
    try:
        data = request.get_json()
        if not data:
            return jsonify({"success": False, "error": "No data provided"}), 400

        if not all(k in data for k in ("name", "email", "password", "role")):
            return jsonify({"success": False, "error": "Missing fields"}), 400
        
        # Check if user already exists
        existing_user = db.get_user(data["email"], data["role"])
        if existing_user:
            return jsonify({"success": False, "error": "Email already exists"}), 400
        
        # Hash password
        hashed_password = bcrypt.hashpw(data["password"].encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
        
        # Create user
        success = db.create_user(
            name=data["name"],
            email=data["email"],
            password=hashed_password,
            role=data["role"]
        )
        
        if success:
            db.log_operation(
                endpoint='/signup',
                operation_type='user_registration',
                user_email=data["email"],
                details=f"User {data['name']} registered as {data['role']}",
                status='success'
            )
            return jsonify({"success": True})
        else:
            return jsonify({"success": False, "error": "Failed to create user"}), 500
    except Exception as e:
        app.logger.error(f"Error in signup endpoint: {str(e)}")
        return jsonify({"success": False, "error": f"Signup failed: {str(e)}"}), 500

@app.route('/signin', methods=['POST'])
def signin():
    """Authenticate a user"""
    try:
        data = request.get_json()
        if not data:
            return jsonify({"success": False, "error": "No data provided"}), 400

        email = data.get("email")
        password = data.get("password")
        role = data.get("role")

        if not email or not password or not role:
            return jsonify({"success": False, "error": "Missing email, password, or role"}), 400

        # Get user from database
        user = db.get_user(email, role)

        if not user:
            # Log attempted login with non-existent user
            db.log_operation(
                endpoint='/signin',
                operation_type='user_login',
                user_email=email,
                details="Attempted login with non-existent user",
                status='failure'
            )
            return jsonify({"success": False, "error": "Invalid email or password"}), 401

        stored_password = user.get("password")
        if stored_password is None:
            db.log_operation(
                endpoint='/signin',
                operation_type='user_login',
                user_email=email,
                details="Account password not set",
                status='error'
            )
            return jsonify({"success": False, "error": "Account password not set"}), 500

        # Support for both str and bytes from DB
        if isinstance(stored_password, str):
            stored_password_bytes = stored_password.encode('utf-8')
        else:
            stored_password_bytes = stored_password

        try:
            # Compare provided password with stored hash
            if bcrypt.checkpw(password.encode('utf-8'), stored_password_bytes):
                db.log_operation(
                    endpoint='/signin',
                    operation_type='user_login',
                    user_email=email,
                    details=f"User {user.get('name', '')} signed in",
                    status='success'
                )
                # It's more standard to return user info or token, but to preserve original logic:
                return jsonify({
                    "success": True,
                    "message": "Login successful",
                    "name": user.get('name', ''),
                    "role": user.get('role', ''),
                    "email": user.get('email', '')
                })
            else:
                db.log_operation(
                    endpoint='/signin',
                    operation_type='user_login',
                    user_email=email,
                    details="Attempted login with incorrect password",
                    status='failure'
                )
                return jsonify({"success": False, "error": "Invalid email or password"}), 401
        except Exception as e:
            app.logger.error(f"Password verification failed: {str(e)}")
            db.log_operation(
                endpoint='/signin',
                operation_type='user_login',
                user_email=email,
                details="Password verification error",
                status='error',
                error_message=str(e)
            )
            return jsonify({"success": False, "error": "Password verification error"}), 500

    except Exception as e:
        app.logger.error(f"Error in signin: {str(e)}")
        db.log_operation(
            endpoint='/signin',
            operation_type='user_login',
            user_email=locals().get("email", None),
            details=f"Unexpected error: {str(e)}",
            status='error',
            error_message=str(e)
        )
        return jsonify({"success": False, "error": "An error occurred during sign in"}), 500

@app.route('/teacher/profile', methods=['POST'])
def create_or_update_teacher_profile():
    """Create or update a teacher profile"""
    try:
        data = request.get_json()
        required_fields = ['user_email', 'name', 'bio']
        if not all(field in data for field in required_fields):
            return jsonify({'success': False, 'error': 'Missing fields'}), 400

        user_email = data['user_email']
        success = db.create_or_update_teacher(
            email=user_email,
            name=data['name'],
            bio=data['bio']
        )

        if success:
            return jsonify({'success': True})
        else:
            return jsonify({'success': False, 'error': 'Failed to save profile'}), 500
            
    except Exception as e:
        app.logger.error(f"Error in create_or_update_teacher_profile: {str(e)}")
        return jsonify({'success': False, 'error': f'An error occurred: {str(e)}'}), 500

@app.route('/teacher/profile/<email>', methods=['GET'])
def get_teacher_profile(email):
    """Get teacher profile by email"""
    try:
        teacher = db.get_teacher(email)
        if not teacher:
            return jsonify({'success': False, 'error': 'Profile not found'}), 404
        
        # Remove id field for response
        profile = {k: v for k, v in teacher.items() if k != 'id'}
        return jsonify({'success': True, 'profile': profile})
        
    except Exception as e:
        app.logger.error(f"Error in get_teacher_profile: {str(e)}")
        return jsonify({'success': False, 'error': f'An error occurred: {str(e)}'}), 500


@app.route('/teacher/materials', methods=['POST'])
def upload_study_material():
    """Upload a study material file"""
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file part'}), 400
        file = request.files['file']
        title = request.form.get('title')
        description = request.form.get('description', '')
        teacher_email = request.form.get('teacher_email')

        if not file or file.filename == '':
            return jsonify({'success': False, 'error': 'No selected file'}), 400
        if not title or not teacher_email:
            return jsonify({'success': False, 'error': 'Missing title or teacher email'}), 400

        filename = secure_filename(file.filename)
        filepath = os.path.join(STUDY_MATERIALS_FOLDER, filename)
        file.save(filepath)

        success = db.add_study_material(
            teacher_email=teacher_email,
            title=title,
            description=description,
            file_path=filepath,
            filename=filename
        )

        if success:
            return jsonify({'success': True})
        else:
            return jsonify({'success': False, 'error': 'Failed to save material'}), 500
            
    except Exception as e:
        app.logger.error(f"Error in upload_study_material: {str(e)}")
        return jsonify({'success': False, 'error': f'An error occurred: {str(e)}'}), 500

@app.route('/teacher/materials/<teacher_email>', methods=['GET'])
def get_study_materials(teacher_email):
    """Get all study materials for a teacher"""
    try:
        materials = db.get_study_materials(teacher_email)
        # Remove id field for response
        materials_list = [{k: v for k, v in mat.items() if k != 'id'} for mat in materials]
        return jsonify({'success': True, 'materials': materials_list})
        
    except Exception as e:
        app.logger.error(f"Error in get_study_materials: {str(e)}")
        return jsonify({'success': False, 'error': f'An error occurred: {str(e)}'}), 500

@app.route('/study_materials/<filename>', methods=['GET'])
def download_study_material(filename):
    return send_from_directory(STUDY_MATERIALS_FOLDER, filename)

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    try:
        # Test database connection
        test_user = db.get_user('test@test.com')
        
        return jsonify({
            'status': 'healthy',
            'database': 'connected',
            'models': {
                'summary': summary_model is not None,
                'question': question_model is not None,
                'quiz': quiz_model is not None
            }
        })
    except Exception as e:
        return jsonify({
            'status': 'unhealthy',
            'error': str(e)
        }), 500

if __name__ == "__main__":
    app.logger.info("=" * 50)
    app.logger.info("Starting Study Buddy Flask Application")
    app.logger.info("=" * 50)
    app.logger.info(f"Server will start on http://{HOST}:{PORT}")
    app.logger.info(f"Debug mode: {FLASK_DEBUG}")
    
    print(f"\n{'='*50}")
    print("Study Buddy AI Summarizer")
    print(f"{'='*50}")
    print(f"Starting server on http://{HOST}:{PORT}")
    print(f"Database: SQLite at {os.path.join(BASE_DIR, 'studybuddy.db')}")
    print(f"Debug mode: {FLASK_DEBUG}")
    print(f"{'='*50}\n")
    
    app.run(debug=FLASK_DEBUG, host=HOST, port=PORT)
